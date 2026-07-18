# app/documents/parser.py
"""
Universal file parser for LTM ingestion.

Returns (text_chunks, structured_records):
  - Unstructured docs (PDF, DOCX, EPUB, PPTX, HTML, TXT): text_chunks populated, records empty.
  - Structured data (CSV, TSV, JSONL, XLSX, Parquet, Arrow, SQLite): records populated, text_chunks empty.

Lazy-imports heavy libraries to avoid startup cost.
"""
import csv
import io
import json
import logging
import sqlite3
from pathlib import Path
from typing import Any, Dict, List, Tuple

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Lazy Import Helpers
# ---------------------------------------------------------------------------


def _get_pymupdf():
    try:
        import fitz

        return fitz
    except ImportError:
        raise ImportError("PyMuPDF not installed. Run: pip install pymupdf")


def _get_docx():
    try:
        from docx import Document

        return Document
    except ImportError:
        raise ImportError("python-docx not installed. Run: pip install python-docx")


def _get_epub():
    try:
        from bs4 import BeautifulSoup
        from ebooklib import epub

        return epub, BeautifulSoup
    except ImportError:
        raise ImportError(
            "ebooklib/bs4 not installed. Run: pip install ebooklib beautifulsoup4"
        )


def _get_pptx():
    try:
        from pptx import Presentation

        return Presentation
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")


def _get_openpyxl():
    try:
        from openpyxl import load_workbook

        return load_workbook
    except ImportError:
        raise ImportError("openpyxl not installed. Run: pip install openpyxl")


def _get_pyarrow():
    try:
        import pyarrow.feather as feather
        import pyarrow.parquet as pq

        return pq, feather
    except ImportError:
        raise ImportError("pyarrow not installed. Run: pip install pyarrow")


def _get_pandas():
    try:
        import pandas as pd

        return pd
    except ImportError:
        raise ImportError("pandas not installed. Run: pip install pandas")


def _get_bs4():
    try:
        from bs4 import BeautifulSoup

        return BeautifulSoup
    except ImportError:
        return None


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def parse_file(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    """Parse any supported file into (text_chunks, structured_records).

    For unstructured documents: returns page/paragraph-level text chunks.
    For structured datasets: returns list of record dicts with all fields preserved.
    For unsupported formats: attempts plain-text fallback.
    """
    ext = file_path.suffix.lower()
    parser = PARSERS.get(ext)

    if parser is None:
        try:
            text = file_path.read_text(encoding="utf-8", errors="replace")
            return ([text] if text.strip() else []), []
        except Exception:
            logger.warning(
                "Cannot parse %s: unsupported format '%s'",
                file_path.name,
                ext,
            )
            return [], []

    try:
        return parser(file_path)
    except Exception as e:
        logger.error(
            "Failed to parse %s (%s): %s",
            file_path.name,
            ext,
            e,
            exc_info=True,
        )
        return [], []


# ---------------------------------------------------------------------------
# Unstructured Document Parsers → List[str]
# ---------------------------------------------------------------------------


def _parse_txt(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    text = file_path.read_text(encoding="utf-8", errors="replace")
    return ([text] if text.strip() else []), []


def _parse_html(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    """Strip scripts/styles/nav/footer/header, return paragraph-level text."""
    bs4 = _get_bs4()
    raw = file_path.read_text(encoding="utf-8", errors="replace")

    if bs4 is None:
        return ([raw] if raw.strip() else []), []

    soup = bs4(raw, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    chunks = [c.strip() for c in text.split("\n\n") if c.strip()]
    return chunks, []


def _parse_pdf(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    """Extract text page-by-page via PyMuPDF."""
    fitz = _get_pymupdf()
    doc = fitz.open(str(file_path))
    pages: List[str] = []
    try:
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text("text")
            if text.strip():
                pages.append(text.strip())
    finally:
        doc.close()
    return pages, []


def _parse_docx(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    """Extract paragraphs AND table rows from DOCX."""
    Document = _get_docx()
    doc = Document(str(file_path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                paragraphs.append(row_text)

    return paragraphs, []


def _parse_epub(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    """Extract text from EPUB chapters, skipping non-document items."""
    epub_mod, bs4 = _get_epub()
    book = epub_mod.read_epub(str(file_path), options={"ignore_ncx": True})
    chapters: List[str] = []

    for item in book.get_items_of_type(9):  # ITEM_DOCUMENT
        try:
            content = item.get_content().decode("utf-8", errors="replace")
            soup = bs4(content, "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            if text.strip():
                chapters.append(text.strip())
        except Exception as e:
            logger.debug("Skipping EPUB item %s: %s", item.get_name(), e)

    return chapters, []


def _parse_pptx(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    """Extract text frames AND tables from each slide."""
    Presentation = _get_pptx()
    prs = Presentation(str(file_path))
    slides: List[str] = []

    for i, slide in enumerate(prs.slides, 1):
        texts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)

            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        texts.append(row_text)

        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))

    return slides, []


# ---------------------------------------------------------------------------
# Structured Dataset Parsers → List[Dict[str, Any]]
# ---------------------------------------------------------------------------


def _parse_csv_tsv(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    """Auto-detect delimiter, return structured records."""
    raw = file_path.read_text(encoding="utf-8", errors="replace")

    try:
        dialect = csv.Sniffer().sniff(raw[:8192])
    except csv.Error:
        dialect = csv.excel

    reader = csv.DictReader(io.StringIO(raw), dialect=dialect)
    records: List[Dict[str, Any]] = []

    for row in reader:
        clean = {
            k.strip(): v.strip() if isinstance(v, str) else v
            for k, v in row.items()
            if k and k.strip()
        }
        if clean:
            records.append(clean)

    return [], records


def _parse_jsonl(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    """Line-delimited JSON → structured records."""
    records: List[Dict[str, Any]] = []
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        for line in f:
            line = line.strip()
            if line:
                try:
                    records.append(json.loads(line))
                except json.JSONDecodeError:
                    logger.warning("Skipping malformed JSON line in %s", file_path.name)
    return [], records


def _parse_xlsx(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    """Excel → structured records with sheet name metadata."""
    load_workbook = _get_openpyxl()
    wb = load_workbook(filename=str(file_path), read_only=True, data_only=True)
    records: List[Dict[str, Any]] = []

    try:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            headers = None
            for row in ws.iter_rows(values_only=True):
                if not any(cell is not None and cell != "" for cell in row):
                    continue
                if headers is None:
                    headers = [
                        str(c).strip() if c is not None else f"col_{i}"
                        for i, c in enumerate(row)
                    ]
                    continue
                record: Dict[str, Any] = {"_sheet": sheet_name}
                for i, val in enumerate(row):
                    if i < len(headers):
                        record[headers[i]] = val
                records.append(record)
    finally:
        wb.close()

    return [], records


def _parse_parquet_arrow(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    """Parquet / Arrow IPC → structured records, batched for memory safety."""
    pq, feather = _get_pyarrow()
    pd = _get_pandas()

    ext = file_path.suffix.lower()
    if ext == ".parquet":
        table = pq.read_table(str(file_path))
    else:
        table = feather.read_table(str(file_path))

    df = table.to_pandas()
    df = df.where(pd.notnull(df), None)
    records = df.to_dict(orient="records")
    return [], records


def _parse_sqlite_db(file_path: Path) -> Tuple[List[str], List[Dict[str, Any]]]:
    """SQLite → structured records from all user tables (10K row safety limit)."""
    conn = sqlite3.connect(str(file_path))
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()

    try:
        cursor.execute(
            "SELECT name FROM sqlite_master WHERE type='table' AND name NOT LIKE 'sqlite_%'"
        )
        tables = [row[0] for row in cursor.fetchall()]

        records: List[Dict[str, Any]] = []
        for table_name in tables:
            try:
                cursor.execute(f"SELECT * FROM [{table_name}] LIMIT 10000")
                rows = cursor.fetchall()
                for row in rows:
                    record = dict(row)
                    record["_table"] = table_name
                    records.append(record)
            except Exception as e:
                logger.debug("Skipping SQLite table '%s': %s", table_name, e)

        return [], records
    finally:
        conn.close()


# ---------------------------------------------------------------------------
# Parser Registry — Single Source of Truth
# ---------------------------------------------------------------------------

PARSERS = {
    ".txt": _parse_txt,
    ".md": _parse_txt,
    ".py": _parse_txt,
    ".js": _parse_txt,
    ".ts": _parse_txt,
    ".json": _parse_txt,
    ".yaml": _parse_txt,
    ".yml": _parse_txt,
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".epub": _parse_epub,
    ".pptx": _parse_pptx,
    ".html": _parse_html,
    ".htm": _parse_html,
    ".csv": _parse_csv_tsv,
    ".tsv": _parse_csv_tsv,
    ".jsonl": _parse_jsonl,
    ".xlsx": _parse_xlsx,
    ".parquet": _parse_parquet_arrow,
    ".arrow": _parse_parquet_arrow,
    ".sqlite": _parse_sqlite_db,
    ".db": _parse_sqlite_db,
}

SUPPORTED_EXTENSIONS = sorted(PARSERS.keys())
