# app/documents/processor.py
import csv
import hashlib
import json
import logging
import os
import re
import xml.etree.ElementTree as ET
from dataclasses import asdict, dataclass
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import pdfplumber
from bs4 import BeautifulSoup
from ebooklib import ITEM_DOCUMENT, epub

try:
    from docx import Document

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    Document = None

try:
    import frontmatter

    FRONTMATTER_AVAILABLE = True
except ImportError:
    FRONTMATTER_AVAILABLE = False
    frontmatter = None

try:
    from string import punctuation

    import spacy
    from spacy.lang.en.stop_words import STOP_WORDS

    SPACY_AVAILABLE = True
    _nlp = None
except ImportError:
    SPACY_AVAILABLE = False
    STOP_WORDS = set()
    punctuation = ""

try:
    from striprtf.striprtf import rtf_to_text

    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False
    rtf_to_text = None

try:
    from odf import opendocument
    from odf.text import P

    ODT_AVAILABLE = True
except ImportError:
    ODT_AVAILABLE = False
    opendocument = None
    P = None

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None

try:
    import openpyxl

    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    openpyxl = None

try:
    from pylatexenc.latex2text import LatexNodes2Text

    LATEX_AVAILABLE = True
except ImportError:
    LATEX_AVAILABLE = False
    LatexNodes2Text = None

try:
    from docutils.core import publish_parts

    RST_AVAILABLE = True
except ImportError:
    RST_AVAILABLE = False
    publish_parts = None

from app.documents.chunker import chunk_text_semantic

logger = logging.getLogger(__name__)


@dataclass
class DocumentMetadata:
    """Rich document metadata structure"""

    filename: str
    source_type: str
    author: Optional[str] = None
    title: Optional[str] = None
    created: Optional[str] = None
    modified: Optional[str] = None
    tags: List[str] = None
    keywords: List[str] = None
    word_count: int = 0
    file_hash: str = ""

    def __post_init__(self):
        if self.tags is None:
            self.tags = []
        if self.keywords is None:
            self.keywords = []

    def to_dict(self) -> Dict:
        """Convert to dictionary for storage"""
        return asdict(self)


def _get_nlp():
    """Lazy-load SpaCy NLP model"""
    global _nlp
    if _nlp is None and SPACY_AVAILABLE:
        try:
            _nlp = spacy.load("en_core_web_sm")
            _nlp.max_length = 2_000_000
            logger.info("SpaCy NLP model loaded successfully")
        except Exception as e:
            logger.warning("Failed to load SpaCy model: %s", e)
            _nlp = None
    return _nlp


def extract_keywords(text: str, top_n: int = 20) -> List[str]:
    """Extract important keywords from text using SpaCy"""
    if not SPACY_AVAILABLE:
        return []

    nlp = _get_nlp()
    if not nlp:
        return []

    try:
        doc = nlp(text.lower())
        candidates = [
            token.text
            for token in doc
            if token.text not in STOP_WORDS
            and token.text not in punctuation
            and len(token.text) > 2
            and token.pos_ in ["NOUN", "PROPN", "ADJ"]
        ]

        freq = {}
        for word in candidates:
            freq[word] = freq.get(word, 0) + 1

        keywords = sorted(freq.items(), key=lambda x: x[1], reverse=True)
        return [kw[0] for kw in keywords[:top_n]]
    except Exception as e:
        logger.warning("Failed to extract keywords: %s", e)
        return []


def clean_text(text: str) -> str:
    """Clean and normalize text, fixing common issues"""
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"([a-z])([A-Z])", r"\1 \2", text)
    return text.strip()


def file_hash(filepath: Path) -> str:
    """Generate SHA256 hash for change detection"""
    h = hashlib.sha256()
    try:
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception as e:
        logger.error("Failed to hash file %s: %s", filepath, e)
        return ""


def extract_metadata(filepath: Path) -> DocumentMetadata:
    """Extract rich metadata from any document format"""
    stat = filepath.stat()
    meta = DocumentMetadata(
        filename=filepath.name,
        source_type=filepath.suffix[1:].lower(),
        created=datetime.fromtimestamp(stat.st_ctime).isoformat(),
        modified=datetime.fromtimestamp(stat.st_mtime).isoformat(),
        file_hash=file_hash(filepath),
    )

    ext = meta.source_type.lower()

    try:
        if ext == "pdf":
            try:
                import fitz

                doc = fitz.open(filepath)
                meta.title = doc.metadata.get("title", "")
                meta.author = doc.metadata.get("author", "")
                doc.close()
            except ImportError:
                pass

        elif ext == "docx" and DOCX_AVAILABLE:
            doc = Document(filepath)
            meta.title = doc.core_properties.title or ""
            meta.author = doc.core_properties.author or ""

        elif ext == "md" and FRONTMATTER_AVAILABLE:
            with open(filepath, "r", encoding="utf-8") as f:
                fm = frontmatter.load(f)
                meta.title = fm.get("title", "")
                meta.author = fm.get("author", "")
                meta.tags = fm.get("tags", [])

        elif ext == "pptx" and PPTX_AVAILABLE:
            prs = Presentation(filepath)
            if prs.core_properties.title:
                meta.title = prs.core_properties.title
            if prs.core_properties.author:
                meta.author = prs.core_properties.author

    except Exception as e:
        logger.warning("Failed to extract metadata from %s: %s", filepath, e)

    return meta


def normalize_text(text: str) -> str:
    """Normalize text by cleaning up whitespace and line endings"""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n\s*\n+", "\n\n", text)
    lines = [line.strip() for line in text.split("\n")]
    return "\n".join(lines).strip()


def process_pdf(pdf_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from PDF and chunk it semantically"""
    chunks = []
    try:
        meta = extract_metadata(pdf_path)
        full_text = ""

        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                try:
                    text = page.extract_text()
                    if text:
                        full_text += text + "\n\n"
                except Exception as e:
                    logger.warning("Skipping page %s in %s: %s", page_num, pdf_path, e)

        full_text = normalize_text(full_text)
        if not full_text:
            logger.warning("No text extracted from PDF %s", pdf_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {
                    "text": chunk_text,
                    "metadata": {
                        **meta.to_dict(),
                        "chunk_index": idx,
                        "chunk_word_count": len(chunk_text.split()),
                    },
                }
            )

        logger.info("PDF '%s' split into %s chunks", pdf_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process PDF %s: %s", pdf_path, e)

    return chunks


def process_epub(epub_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from EPUB and chunk it semantically"""
    chunks = []
    try:
        meta = extract_metadata(epub_path)
        book = epub.read_epub(epub_path)
        full_text = ""
        section_titles = []

        for item in book.get_items():
            if item.get_type() == ITEM_DOCUMENT:
                try:
                    html_content = item.get_content().decode("utf-8")
                    soup = BeautifulSoup(html_content, "html.parser")
                    headings = [
                        h.get_text(strip=True)
                        for h in soup.find_all(["h1", "h2", "h3"])
                    ]
                    if headings:
                        section_titles.extend(headings)

                    paragraphs = []
                    for tag in soup.find_all(["p", "div", "section", "article"]):
                        text = tag.get_text(separator=" ", strip=True)
                        if text:
                            paragraphs.append(text)
                    full_text += "\n\n".join(paragraphs) + "\n\n"
                except Exception as e:
                    logger.warning("Skipping item in EPUB %s: %s", epub_path, e)

        full_text = normalize_text(full_text)
        if not full_text:
            logger.warning("No text extracted from EPUB %s", epub_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            section_title = (
                section_titles[idx - 1] if idx - 1 < len(section_titles) else None
            )
            chunks.append(
                {
                    "text": chunk_text,
                    "metadata": {
                        **meta.to_dict(),
                        "chunk_index": idx,
                        "section_title": section_title,
                    },
                }
            )

        logger.info("EPUB '%s' split into %s chunks", epub_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process EPUB %s: %s", epub_path, e)

    return chunks


def process_txt(txt_path: Path, max_words: int = 300) -> List[Dict]:
    """Read TXT file and chunk it semantically"""
    chunks = []
    try:
        meta = extract_metadata(txt_path)

        try:
            with open(txt_path, "r", encoding="utf-8") as f:
                full_text = f.read()
        except UnicodeDecodeError:
            with open(txt_path, "r", encoding="cp1252") as f:
                full_text = f.read()

        full_text = normalize_text(full_text)
        if not full_text:
            logger.warning("No text extracted from TXT %s", txt_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("TXT '%s' split into %s chunks", txt_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process TXT %s: %s", txt_path, e)

    return chunks


def process_docx(docx_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from DOCX and chunk it semantically"""
    if not DOCX_AVAILABLE:
        logger.warning("python-docx not installed, skipping DOCX processing")
        return []

    chunks = []
    try:
        meta = extract_metadata(docx_path)
        doc = Document(docx_path)

        full_text = "\n\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        full_text = normalize_text(full_text)

        if not full_text:
            logger.warning("No text extracted from DOCX %s", docx_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("DOCX '%s' split into %s chunks", docx_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process DOCX %s: %s", docx_path, e)

    return chunks


def process_html(html_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from HTML and chunk it semantically"""
    chunks = []
    try:
        meta = extract_metadata(html_path)

        with open(html_path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")

            if not meta.title and soup.title:
                meta.title = soup.title.get_text(strip=True)

            body = soup.find("body")
            if body:
                full_text = body.get_text(separator="\n\n", strip=True)
            else:
                full_text = soup.get_text(separator="\n\n", strip=True)

        full_text = normalize_text(full_text)
        if not full_text:
            logger.warning("No text extracted from HTML %s", html_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("HTML '%s' split into %s chunks", html_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process HTML %s: %s", html_path, e)

    return chunks


def process_markdown(md_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from Markdown (with optional frontmatter) and chunk it"""
    chunks = []
    try:
        meta = extract_metadata(md_path)

        with open(md_path, "r", encoding="utf-8") as f:
            if FRONTMATTER_AVAILABLE:
                post = frontmatter.load(f)
                full_text = post.content
            else:
                full_text = f.read()

        full_text = normalize_text(full_text)
        if not full_text:
            logger.warning("No text extracted from Markdown %s", md_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("Markdown '%s' split into %s chunks", md_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process Markdown %s: %s", md_path, e)

    return chunks


def process_rtf(rtf_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from RTF and chunk it semantically"""
    if not RTF_AVAILABLE:
        logger.warning("striprtf not installed, skipping RTF processing")
        return []

    chunks = []
    try:
        meta = extract_metadata(rtf_path)

        with open(rtf_path, "r", encoding="utf-8", errors="ignore") as f:
            rtf_content = f.read()

        full_text = rtf_to_text(rtf_content)
        full_text = normalize_text(full_text)

        if not full_text:
            logger.warning("No text extracted from RTF %s", rtf_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("RTF '%s' split into %s chunks", rtf_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process RTF %s: %s", rtf_path, e)

    return chunks


def process_odt(odt_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from ODT (OpenDocument) and chunk it semantically"""
    if not ODT_AVAILABLE:
        logger.warning("odfpy not installed, skipping ODT processing")
        return []

    chunks = []
    try:
        meta = extract_metadata(odt_path)

        doc = opendocument.load(str(odt_path))
        full_text = "\n\n".join(
            [p.text for p in doc.text.getElementsByType(P) if p.text.strip()]
        )
        full_text = normalize_text(full_text)

        if not full_text:
            logger.warning("No text extracted from ODT %s", odt_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("ODT '%s' split into %s chunks", odt_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process ODT %s: %s", odt_path, e)

    return chunks


def process_csv(csv_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from CSV/TSV and chunk it semantically"""
    chunks = []
    try:
        meta = extract_metadata(csv_path)

        with open(csv_path, "r", encoding="utf-8", errors="ignore") as f:
            sample = f.read(2048)
            try:
                dialect = csv.Sniffer().sniff(sample)
                delimiter = dialect.delimiter
            except csv.Error:
                delimiter = ","

        with open(csv_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f, delimiter=delimiter)
            rows = list(reader)

        if not rows:
            logger.warning("No data in CSV %s", csv_path)
            return chunks

        full_text = ""
        headers = rows[0] if rows else []
        full_text += "Headers: " + ", ".join(headers) + "\n\n"

        for row in rows[1:]:
            row_text = " | ".join(row)
            full_text += row_text + "\n"

        full_text = normalize_text(full_text)
        if not full_text:
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("CSV '%s' split into %s chunks", csv_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process CSV %s: %s", csv_path, e)

    return chunks


def process_json(json_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from JSON and chunk it semantically"""
    chunks = []
    try:
        meta = extract_metadata(json_path)

        with open(json_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        def json_to_text(obj, indent=0):
            text = ""
            if isinstance(obj, dict):
                for key, value in obj.items():
                    text += f"{'  ' * indent}{key}: "
                    if isinstance(value, (dict, list)):
                        text += "\n" + json_to_text(value, indent + 1)
                    else:
                        text += f"{value}\n"
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    text += f"{'  ' * indent}[{i}] "
                    if isinstance(item, (dict, list)):
                        text += "\n" + json_to_text(item, indent + 1)
                    else:
                        text += f"{item}\n"
            else:
                text += f"{'  ' * indent}{obj}\n"
            return text

        full_text = json_to_text(data)
        full_text = normalize_text(full_text)

        if not full_text:
            logger.warning("No text extracted from JSON %s", json_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("JSON '%s' split into %s chunks", json_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process JSON %s: %s", json_path, e)

    return chunks


def process_xml(xml_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from XML and chunk it semantically"""
    chunks = []
    try:
        meta = extract_metadata(xml_path)

        tree = ET.parse(xml_path)
        root = tree.getroot()

        full_text = ""
        for elem in root.iter():
            if elem.text and elem.text.strip():
                full_text += elem.text.strip() + "\n"
            if elem.tail and elem.tail.strip():
                full_text += elem.tail.strip() + "\n"

        full_text = normalize_text(full_text)

        if not full_text:
            logger.warning("No text extracted from XML %s", xml_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("XML '%s' split into %s chunks", xml_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process XML %s: %s", xml_path, e)

    return chunks


def process_pptx(pptx_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from PowerPoint and chunk it semantically"""
    if not PPTX_AVAILABLE:
        logger.warning("python-pptx not installed, skipping PPTX processing")
        return []

    chunks = []
    try:
        meta = extract_metadata(pptx_path)
        prs = Presentation(pptx_path)

        full_text = ""
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = f"Slide {slide_num}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text.strip() + "\n"
            full_text += slide_text + "\n\n"

        full_text = normalize_text(full_text)

        if not full_text:
            logger.warning("No text extracted from PPTX %s", pptx_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("PPTX '%s' split into %s chunks", pptx_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process PPTX %s: %s", pptx_path, e)

    return chunks


def process_xlsx(xlsx_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from Excel and chunk it semantically"""
    if not XLSX_AVAILABLE:
        logger.warning("openpyxl not installed, skipping XLSX processing")
        return []

    chunks = []
    try:
        meta = extract_metadata(xlsx_path)
        wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)

        full_text = ""
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            full_text += f"Sheet: {sheet_name}\n"

            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(
                    [str(cell) if cell is not None else "" for cell in row]
                )
                if row_text.strip("| "):
                    full_text += row_text + "\n"
            full_text += "\n"

        wb.close()
        full_text = normalize_text(full_text)

        if not full_text:
            logger.warning("No text extracted from XLSX %s", xlsx_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("XLSX '%s' split into %s chunks", xlsx_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process XLSX %s: %s", xlsx_path, e)

    return chunks


def process_tex(tex_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from LaTeX and chunk it semantically"""
    chunks = []
    try:
        meta = extract_metadata(tex_path)

        with open(tex_path, "r", encoding="utf-8", errors="ignore") as f:
            latex_content = f.read()

        if LATEX_AVAILABLE:
            full_text = LatexNodes2Text().latex_to_text(latex_content)
        else:
            full_text = re.sub(r"\\[a-zA-Z]+\{[^}]*\}", "", latex_content)
            full_text = re.sub(r"\\[a-zA-Z]+", "", full_text)
            full_text = re.sub(r"[{}]", "", full_text)
            full_text = re.sub(r"%%.*?\n", "\n", full_text)

        full_text = normalize_text(full_text)

        if not full_text:
            logger.warning("No text extracted from LaTeX %s", tex_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("LaTeX '%s' split into %s chunks", tex_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process LaTeX %s: %s", tex_path, e)

    return chunks


def process_rst(rst_path: Path, max_words: int = 300) -> List[Dict]:
    """Extract text from reStructuredText and chunk it semantically"""
    chunks = []
    try:
        meta = extract_metadata(rst_path)

        with open(rst_path, "r", encoding="utf-8") as f:
            rst_content = f.read()

        if RST_AVAILABLE:
            parts = publish_parts(rst_content, writer_name="html")
            soup = BeautifulSoup(parts["body"], "html.parser")
            full_text = soup.get_text()
        else:
            full_text = re.sub(
                r"^[=\-~`\"+#*^]{3,}\s*$", "", rst_content, flags=re.MULTILINE
            )
            full_text = re.sub(r"\.\.\s+[a-z]+::", "", full_text)
            full_text = re.sub(r":[a-z]+:`[^`]*`", "", full_text)

        full_text = normalize_text(full_text)

        if not full_text:
            logger.warning("No text extracted from RST %s", rst_path)
            return chunks

        meta.word_count = len(full_text.split())
        meta.keywords = extract_keywords(full_text[:5000])

        chunk_texts = chunk_text_semantic(
            full_text, max_words=max_words, overlap_sentences=2
        )

        for idx, chunk_text in enumerate(chunk_texts, start=1):
            chunks.append(
                {"text": chunk_text, "metadata": {**meta.to_dict(), "chunk_index": idx}}
            )

        logger.info("RST '%s' split into %s chunks", rst_path.name, len(chunks))
    except Exception as e:
        logger.error("Failed to process RST %s: %s", rst_path, e)

    return chunks


def process_directory(
    directory: Path, max_words: int = 300, force_reprocess: bool = False
) -> List[Dict]:
    """Process all supported files in a directory recursively with change detection"""
    all_chunks = []
    processed_count = 0
    skipped_count = 0

    metadata_cache_dir = directory / ".metadata_cache"
    metadata_cache_dir.mkdir(parents=True, exist_ok=True)

    for file_path in directory.rglob("*"):
        if not file_path.is_file():
            continue

        if file_path.name.startswith(".") or ".metadata_cache" in str(file_path):
            continue

        ext = file_path.suffix.lower()

        if not force_reprocess:
            cache_file = metadata_cache_dir / f"{file_path.name}.json"
            if cache_file.exists():
                try:
                    with open(cache_file, "r") as f:
                        cached = json.load(f)
                    current_hash = file_hash(file_path)
                    if cached.get("file_hash") == current_hash:
                        logger.debug("Skipping unchanged file: %s", file_path.name)
                        skipped_count += 1
                        continue
                except Exception as e:
                    logger.warning("Failed to read cache for %s: %s", file_path, e)

        file_chunks = []
        if ext == ".pdf":
            file_chunks = process_pdf(file_path, max_words)
        elif ext == ".epub":
            file_chunks = process_epub(file_path, max_words)
        elif ext == ".txt":
            file_chunks = process_txt(file_path, max_words)
        elif ext == ".docx":
            file_chunks = process_docx(file_path, max_words)
        elif ext in [".html", ".htm"]:
            file_chunks = process_html(file_path, max_words)
        elif ext == ".md":
            file_chunks = process_markdown(file_path, max_words)
        elif ext == ".rtf":
            file_chunks = process_rtf(file_path, max_words)
        elif ext == ".odt":
            file_chunks = process_odt(file_path, max_words)
        elif ext in [".csv", ".tsv"]:
            file_chunks = process_csv(file_path, max_words)
        elif ext == ".json":
            file_chunks = process_json(file_path, max_words)
        elif ext == ".xml":
            file_chunks = process_xml(file_path, max_words)
        elif ext == ".pptx":
            file_chunks = process_pptx(file_path, max_words)
        elif ext == ".xlsx":
            file_chunks = process_xlsx(file_path, max_words)
        elif ext == ".tex":
            file_chunks = process_tex(file_path, max_words)
        elif ext == ".rst":
            file_chunks = process_rst(file_path, max_words)

        if file_chunks:
            all_chunks.extend(file_chunks)
            processed_count += 1

            try:
                cache_file = metadata_cache_dir / f"{file_path.name}.json"
                with open(cache_file, "w") as f:
                    json.dump(
                        {
                            "file_hash": file_hash(file_path),
                            "processed_at": datetime.now().isoformat(),
                            "chunk_count": len(file_chunks),
                        },
                        f,
                    )
            except Exception as e:
                logger.warning("Failed to update cache for %s: %s", file_path, e)

    logger.info(
        "Processed %s files (%s unchanged) → %s chunks from %s",
        processed_count,
        skipped_count,
        len(all_chunks),
        directory,
    )
    return all_chunks


def get_supported_extensions() -> List[str]:
    """Return list of supported file extensions"""
    return [
        ".pdf",
        ".epub",
        ".txt",
        ".docx",
        ".html",
        ".htm",
        ".md",
        ".rtf",
        ".odt",
        ".csv",
        ".tsv",
        ".json",
        ".xml",
        ".pptx",
        ".xlsx",
        ".tex",
        ".rst",
    ]
